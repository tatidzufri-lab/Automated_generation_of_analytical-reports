from __future__ import annotations

import logging
from typing import Any

from pptx import Presentation
from pptx.util import Inches, Pt

from app.core.config import Settings, get_settings
from app.services.artifact_naming import humanize_artifact_name
from app.services.file_service import FileService, StoredFile
from app.services.time_utils import filename_timestamp, format_local


logger = logging.getLogger(__name__)


class PptxServiceError(RuntimeError):
    """Raised when PPTX generation fails."""


class PptxService:
    """Python-pptx based presentation generator for business analytics."""

    def __init__(self, file_service: FileService, settings: Settings | None = None) -> None:
        self.file_service = file_service
        self.settings = settings or get_settings()

    def generate_report(
        self,
        stored_file: StoredFile,
        analysis: dict[str, Any],
        business_metrics: dict[str, Any] | None,
        charts: list[dict[str, Any]],
        title: str | None = None,
    ) -> dict[str, str]:
        self.file_service.ensure_storage()
        report_name = self._build_report_name(stored_file.file_id)
        output_path = self.settings.output_dir / report_name

        title = title or f"Аналитический отчёт · {stored_file.original_name}"

        try:
            prs = Presentation()
            self._add_title_slide(prs, title)
            self._add_metrics_slide(prs, analysis, business_metrics)

            for chart in charts:
                self._add_chart_slide(prs, chart)

            top_items = (business_metrics or {}).get("top_items") or []
            if top_items:
                self._add_top_items_table_slide(prs, top_items)

            insights = analysis.get("insights", []) if analysis else []
            if insights:
                self._add_insights_slide(prs, insights)

            prs.save(str(output_path))
        except Exception as exc:  # pragma: no cover - external tool errors
            logger.exception("Failed to generate PPTX for %s", stored_file.file_id)
            raise PptxServiceError(f"Не удалось сгенерировать PPTX: {exc}") from exc

        logger.info("Generated PPTX %s for %s", report_name, stored_file.file_id)
        return {
            "title": "PowerPoint deck",
            "kind": "pptx",
            "file_name": report_name,
            "display_name": humanize_artifact_name(report_name),
            "description": "Презентация с метриками, графиками и таблицей топ позиций.",
            "storage_url": f"/storage/outputs/{report_name}",
            "download_url": f"/download/{report_name}",
        }

    # ------------------------------------------------------------------
    # Slide builders
    # ------------------------------------------------------------------
    def _add_title_slide(self, prs: Presentation, title: str) -> None:
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = title
        subtitle = slide.placeholders[1]
        subtitle.text = f"Сгенерирован: {format_local('%d.%m.%Y %H:%M')}"

    def _add_metrics_slide(
        self,
        prs: Presentation,
        analysis: dict[str, Any],
        business_metrics: dict[str, Any] | None,
    ) -> None:
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "Ключевые метрики"
        body = slide.placeholders[1].text_frame
        body.clear()

        lines: list[str] = []
        if business_metrics and (
            business_metrics.get("total_sales") or business_metrics.get("total_orders")
        ):
            lines.append(f"• Общая сумма: {business_metrics.get('total_sales', 0):,.2f}".replace(",", " "))
            lines.append(f"• Средний чек: {business_metrics.get('avg_ticket', 0):,.2f}".replace(",", " "))
            lines.append(f"• Количество записей: {business_metrics.get('total_orders', 0):,}".replace(",", " "))

        summary = analysis.get("summary") if analysis else None
        if summary:
            lines.append(f"• Строк / колонок: {summary.get('rows', 0)} × {summary.get('columns', 0)}")
            lines.append(f"• Числовых метрик: {summary.get('numeric_columns', 0)}")
            lines.append(f"• Пропусков: {summary.get('missing_cells', 0)}")

        if not lines:
            lines.append("• Недостаточно данных для сводных метрик.")

        first = body.paragraphs[0]
        first.text = lines[0]
        first.font.size = Pt(18)
        for extra in lines[1:]:
            paragraph = body.add_paragraph()
            paragraph.text = extra
            paragraph.font.size = Pt(18)

    def _add_chart_slide(self, prs: Presentation, chart: dict[str, Any]) -> None:
        image_path = self.settings.output_dir / chart.get("file_name", "")
        if not image_path.exists():
            return

        slide = prs.slides.add_slide(prs.slide_layouts[5])
        self._remove_placeholders(slide)

        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.9))
        title_frame = title_box.text_frame
        title_frame.text = (
            chart.get("display_name")
            or chart.get("title")
            or humanize_artifact_name(chart.get("file_name", "Chart"))
        )
        title_frame.paragraphs[0].font.size = Pt(24)
        title_frame.paragraphs[0].font.bold = True

        slide.shapes.add_picture(
            str(image_path),
            Inches(0.5),
            Inches(1.3),
            width=Inches(9),
            height=Inches(5.2),
        )

        caption = chart.get("description")
        if caption:
            caption_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.6), Inches(9), Inches(0.6))
            caption_frame = caption_box.text_frame
            caption_frame.text = caption
            caption_frame.paragraphs[0].font.size = Pt(12)
            caption_frame.paragraphs[0].font.italic = True

    def _add_top_items_table_slide(self, prs: Presentation, top_items: list[dict[str, Any]]) -> None:
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        self._remove_placeholders(slide)

        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.9))
        title_frame = title_box.text_frame
        title_frame.text = "Топ позиций"
        title_frame.paragraphs[0].font.size = Pt(24)
        title_frame.paragraphs[0].font.bold = True

        rows = len(top_items) + 1
        table = slide.shapes.add_table(rows, 2, Inches(1), Inches(1.4), Inches(8), Inches(5.2)).table
        table.cell(0, 0).text = "Позиция"
        table.cell(0, 1).text = "Сумма"
        for col_index in range(2):
            table.cell(0, col_index).text_frame.paragraphs[0].font.bold = True

        for index, item in enumerate(top_items, start=1):
            table.cell(index, 0).text = str(item.get("item", ""))
            amount = float(item.get("amount", 0))
            table.cell(index, 1).text = f"{amount:,.2f}".replace(",", " ")

    def _add_insights_slide(self, prs: Presentation, insights: list[str]) -> None:
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "Выводы"
        body = slide.placeholders[1].text_frame
        body.clear()
        body.paragraphs[0].text = f"• {insights[0]}"
        body.paragraphs[0].font.size = Pt(16)
        for insight in insights[1:]:
            paragraph = body.add_paragraph()
            paragraph.text = f"• {insight}"
            paragraph.font.size = Pt(16)

    def _remove_placeholders(self, slide) -> None:
        sp_tree = slide.shapes._spTree
        for shape in list(slide.shapes):
            if shape.is_placeholder:
                sp_tree.remove(shape._element)

    def _build_report_name(self, file_id: str) -> str:
        return f"{file_id}__presentation__{filename_timestamp()}.pptx"
