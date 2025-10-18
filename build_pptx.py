"""
Модуль для генерации PowerPoint презентаций.
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os
from typing import Dict


def build_pptx(context: Dict, output_pptx: str) -> None:
    """
    Создаёт PowerPoint презентацию с аналитическим отчётом.
    
    Args:
        context: Словарь с данными для презентации
        output_pptx: Путь для сохранения PPTX файла
        
    Raises:
        Exception: При ошибках генерации PPTX
    """
    try:
        # Создаём новую презентацию
        prs = Presentation()
        
        # 1. Титульный слайд
        _add_title_slide(prs, context)
        
        # 2. Слайд с ключевыми метриками
        _add_metrics_slide(prs, context)
        
        # 3. Слайд с динамикой продаж (если есть график)
        if context.get('timeseries_png') and os.path.exists(context['timeseries_png']):
            _add_timeseries_slide(prs, context)
        
        # 4. Слайд с топ позициями
        if context.get('top_items_png') and os.path.exists(context['top_items_png']):
            _add_top_items_slide(prs, context)
        
        # 5. Слайд с таблицей топ позиций
        top_items = context.get('top_items')
        if top_items is not None and len(top_items) > 0:
            _add_top_items_table_slide(prs, context)
        
        # Создаём директорию для выходного файла если её нет
        os.makedirs(os.path.dirname(output_pptx), exist_ok=True)
        
        # Сохраняем презентацию
        prs.save(output_pptx)
        
        print(f"PowerPoint презентация сохранена: {output_pptx}")
        
    except Exception as e:
        raise Exception(f"Ошибка при создании PPTX: {str(e)}")


def _add_title_slide(prs: Presentation, context: Dict) -> None:
    """Добавляет титульный слайд."""
    slide_layout = prs.slide_layouts[0]  # Титульный макет
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = context.get('title', 'Аналитический отчёт')
    subtitle.text = f"Сгенерирован: {context.get('generated_at', 'Неизвестно')}"


def _add_metrics_slide(prs: Presentation, context: Dict) -> None:
    """Добавляет слайд с ключевыми метриками."""
    slide_layout = prs.slide_layouts[1]  # Заголовок и содержимое
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    title.text = "Ключевые метрики"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.clear()
    
    # Форматируем числа
    total_sales = context.get('total_sales', 0)
    avg_ticket = context.get('avg_ticket', 0)
    total_orders = context.get('total_orders', 0)
    
    metrics_text = f"""• Общая сумма продаж: {total_sales:,.2f} руб.
• Средний чек: {avg_ticket:,.2f} руб.
• Общее количество заказов: {total_orders:,}"""
    
    p = tf.paragraphs[0]
    p.text = metrics_text
    p.font.size = Pt(18)


def _add_timeseries_slide(prs: Presentation, context: Dict) -> None:
    """Добавляет слайд с графиком динамики продаж."""
    slide_layout = prs.slide_layouts[5]  # Пустой слайд
    slide = prs.slides.add_slide(slide_layout)
    
    # Добавляем заголовок
    title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
    title_frame = title_shape.text_frame
    title_frame.text = "Динамика продаж"
    title_frame.paragraphs[0].font.size = Pt(24)
    title_frame.paragraphs[0].font.bold = True
    
    # Добавляем изображение
    img_path = context['timeseries_png']
    if os.path.exists(img_path):
        slide.shapes.add_picture(
            img_path, 
            Inches(0.5), Inches(1.5), 
            width=Inches(9), 
            height=Inches(5)
        )


def _add_top_items_slide(prs: Presentation, context: Dict) -> None:
    """Добавляет слайд с графиком топ позиций."""
    slide_layout = prs.slide_layouts[5]  # Пустой слайд
    slide = prs.slides.add_slide(slide_layout)
    
    # Добавляем заголовок
    title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
    title_frame = title_shape.text_frame
    title_frame.text = "Топ позиций по продажам"
    title_frame.paragraphs[0].font.size = Pt(24)
    title_frame.paragraphs[0].font.bold = True
    
    # Добавляем изображение
    img_path = context['top_items_png']
    if os.path.exists(img_path):
        slide.shapes.add_picture(
            img_path, 
            Inches(0.5), Inches(1.5), 
            width=Inches(9), 
            height=Inches(5)
        )


def _add_top_items_table_slide(prs: Presentation, context: Dict) -> None:
    """Добавляет слайды с таблицей топ позиций (разбивает на 2 слайда)."""
    top_items = context.get('top_items')
    if top_items is not None and len(top_items) > 0:
        # Разбиваем данные на две части
        mid_point = len(top_items) // 2
        first_half = top_items[:mid_point]
        second_half = top_items[mid_point:]
        
        # Первый слайд
        slide_layout = prs.slide_layouts[5]  # Пустой слайд
        slide1 = prs.slides.add_slide(slide_layout)
        
        # Заголовок первого слайда
        title_shape1 = slide1.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
        title_frame1 = title_shape1.text_frame
        title_frame1.text = "Детализация топ позиций (часть 1)"
        title_frame1.paragraphs[0].font.size = Pt(24)
        title_frame1.paragraphs[0].font.bold = True
        
        # Таблица первого слайда
        rows1 = len(first_half) + 1  # +1 для заголовка
        cols = 2
        
        table1 = slide1.shapes.add_table(
            rows1, cols, 
            Inches(1), Inches(1.5), 
            Inches(8), Inches(4.5)
        ).table
        
        # Заголовки
        table1.cell(0, 0).text = "Позиция"
        table1.cell(0, 1).text = "Сумма продаж"
        
        # Данные первого слайда
        for i, item in enumerate(first_half, 1):
            table1.cell(i, 0).text = str(item['item'])
            table1.cell(i, 1).text = f"{item['amount']:,.2f}"
        
        # Форматирование заголовков первого слайда
        for j in range(cols):
            table1.cell(0, j).text_frame.paragraphs[0].font.bold = True
        
        # Второй слайд (если есть вторая половина)
        if second_half:
            slide2 = prs.slides.add_slide(slide_layout)
            
            # Заголовок второго слайда
            title_shape2 = slide2.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
            title_frame2 = title_shape2.text_frame
            title_frame2.text = "Детализация топ позиций (часть 2)"
            title_frame2.paragraphs[0].font.size = Pt(24)
            title_frame2.paragraphs[0].font.bold = True
            
            # Таблица второго слайда
            rows2 = len(second_half) + 1  # +1 для заголовка
            
            table2 = slide2.shapes.add_table(
                rows2, cols, 
                Inches(1), Inches(1.5), 
                Inches(8), Inches(4.5)
            ).table
            
            # Заголовки
            table2.cell(0, 0).text = "Позиция"
            table2.cell(0, 1).text = "Сумма продаж"
            
            # Данные второго слайда
            for i, item in enumerate(second_half, 1):
                table2.cell(i, 0).text = str(item['item'])
                table2.cell(i, 1).text = f"{item['amount']:,.2f}"
            
            # Форматирование заголовков второго слайда
            for j in range(cols):
                table2.cell(0, j).text_frame.paragraphs[0].font.bold = True
